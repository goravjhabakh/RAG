from pptx import Presentation
from converters.base import BaseDocumentConverter
from typing import Dict

class PPTToMarkdownConverter(BaseDocumentConverter):
    def convert(self, file_path: str) -> Dict[str, str]:
        prs = Presentation(file_path)
        slides_md = []

        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text.strip())
            slides_md.append(f"### Slide {idx}\n" + "\n".join(texts))

        return {"slides": "\n\n".join(slides_md)}